#!/usr/bin/env python3
"""
文件内容提取器（解读模式）
支持: PDF, DOCX, XLSX, PPTX, TXT
用法: python scripts/file_reader.py <fileId>
输出提取的文字内容到 stdout（纯文本）
"""
import os, sys, json

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

UPLOAD_DIR = os.path.expanduser("~/.openclaw/uploads")
WORKSPACE = os.path.expanduser("~/.openclaw/workspace")


def resolve_file(file_id: str):
    for d in [UPLOAD_DIR, os.path.join(WORKSPACE, "storage")]:
        if not os.path.isdir(d):
            continue
        for fn in os.listdir(d):
            if fn.startswith(file_id):
                return os.path.join(d, fn), fn
    return None, None


def read_pdf(path: str) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"--- 第 {i+1} 页 ---\n{text.strip()}")
    return "\n\n".join(pages) if pages else "(PDF 无可提取文字，可能是扫描件)"


def read_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    # 也读表格
    for table in doc.tables:
        paras.append("---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paras.append(" | ".join(cells))
    return "\n".join(paras) if paras else "(空文档)"


def read_xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    lines = []
    for name in wb.sheetnames:
        ws = wb[name]
        lines.append(f"=== Sheet: {name} ===")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c for c in cells):
                lines.append("\t".join(cells))
        lines.append("")
    return "\n".join(lines) if lines else "(空表格)"


def read_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(cells))
        if texts:
            slides.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "(无文字内容)"


def read_txt(path: str) -> str:
    for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(path, 'r', encoding=enc) as f:
                return f.read(10000)
        except:
            continue
    return "(无法读取)"


# ── 主入口 ──
if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({"error": "缺少 fileId"}, ensure_ascii=False))
        sys.exit(1)

    file_id = sys.argv[1]
    path, filename = resolve_file(file_id)
    if not path:
        print(json.dumps({"error": f"未找到文件: {file_id}"}, ensure_ascii=False))
        sys.exit(1)

    ext = os.path.splitext(filename)[1].lower()
    readers = {
        '.pdf': read_pdf,
        '.docx': read_docx,
        '.xlsx': read_xlsx,
        '.pptx': read_pptx,
        '.txt': read_txt,
        '.md': read_txt,
        '.csv': read_txt,
    }

    reader = readers.get(ext)
    if not reader:
        print(json.dumps({"error": f"不支持解析: {ext}", "filename": filename}, ensure_ascii=False))
        sys.exit(1)

    try:
        text = reader(path)
        result = {
            "ok": True,
            "filename": filename,
            "type": ext,
            "chars": len(text),
            "content": text
        }
        print(json.dumps(result, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({"error": f"解析失败: {e}", "filename": filename}, ensure_ascii=False))
