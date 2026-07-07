#!/usr/bin/env python3
"""
文件查看器 — PDF / TXT / 图片 白板展示
用法: python scripts/file_viewer.py <fileId>
自动检测格式，生成对应 viewer HTML，输出 MODAL url。
"""
import os, sys, json, shutil, re, urllib.parse
from pathlib import Path

sys.stdout.reconfigure(encoding='utf-8')

UPLOAD_DIR = os.path.expanduser("~/.openclaw/uploads")
CANVAS_DIR = os.path.join(
    os.environ.get("OPENCLAW_STATE_DIR", os.path.expanduser("~/.openclaw")),
    "canvas",
)
WORKSPACE = os.path.expanduser("~/.openclaw/workspace")

def cleanup_old_viewers():
    """清理旧的 viewer HTML 文件"""
    if not os.path.isdir(CANVAS_DIR):
        return
    for fn in os.listdir(CANVAS_DIR):
        if fn.startswith('viewer_') and fn.endswith('.html'):
            os.remove(os.path.join(CANVAS_DIR, fn))


def resolve_file(file_id: str) -> tuple[str, str] | None:
    """在 uploads 中按 fileId 前缀查找文件"""
    for d in [UPLOAD_DIR, os.path.join(WORKSPACE, "storage")]:
        if not os.path.isdir(d):
            continue
        for fn in os.listdir(d):
            if fn.startswith(file_id):
                return os.path.join(d, fn), fn
    # 模糊搜索
    for d in [UPLOAD_DIR, os.path.join(WORKSPACE, "storage")]:
        if not os.path.isdir(d):
            continue
        for fn in os.listdir(d):
            if file_id.lower() in fn.lower():
                return os.path.join(d, fn), fn
    return None


def copy_to_canvas(src: str, filename: str) -> str:
    dst = os.path.join(CANVAS_DIR, filename)
    if not os.path.exists(dst):
        shutil.copy2(src, dst)
    return dst


def viewer_txt(filepath: str, filename: str) -> str:
    content = ""
    for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
        try:
            with open(filepath, 'r', encoding=enc) as f:
                content = f.read(50000)
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    escaped = content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    return f'''<!DOCTYPE html><html lang="zh"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1,maximum-scale=1,user-scalable=no">
<title>{filename}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;color:#d2d8e0;padding:16px;font-family:'SF Mono',Consolas,'Microsoft YaHei',monospace;font-size:14px;line-height:1.7;white-space:pre-wrap;word-break:break-all;min-height:100dvh}}
.header{{color:#8b949e;font-size:12px;margin-bottom:16px;padding-bottom:8px;border-bottom:1px solid #21262d}}
</style></head><body>
<div class="header">📄 {filename} · {len(content)} 字符</div>
{escaped}
</body></html>'''


def viewer_pdf(filepath: str, filename: str, file_id: str = '', title: str = '', port: int = 18889) -> str:
    # 复制到 canvas（供 pdf.js 渲染）+ 复制到 storage（供 upload_server 下载）
    safe_name = f'{file_id}.pdf'
    copy_to_canvas(filepath, safe_name)
    # 放到 storage 供 upload_server 下载，并写入 index.json（标签=临时）
    storage_dir = os.path.join(WORKSPACE, 'storage')
    storage_path = os.path.join(storage_dir, safe_name)
    if not os.path.exists(storage_path):
        shutil.copy2(filepath, storage_path)
    # 写 index.json
    idx_file = os.path.join(storage_dir, 'index.json')
    idx = {}
    if os.path.exists(idx_file):
        with open(idx_file, 'r', encoding='utf-8') as f:
            idx = json.load(f)
    items = idx.get('items', [])
    items = [i for i in items if i.get('fileName') != safe_name]
    from datetime import datetime
    items.append({'fileId': file_id, 'fileName': safe_name, 'displayName': filename,
                  'savedAt': datetime.now().isoformat(), 'path': storage_path, 'tags': ['临时']})
    idx['items'] = items
    with open(idx_file, 'w', encoding='utf-8') as f:
        json.dump(idx, f, ensure_ascii=False, indent=2)
    pdf_url = f'/__openclaw__/canvas/{safe_name}'
    # 下载走 upload_server 18889（和画廊/游记 PDF 一致）
    download_url = f'http://__HOST__:{port}/file/download/{safe_name}'
    # displayName 优先用传入的 title（用户友好），fallback 到安全文件名
    disp = f'{title}.pdf' if title and title != file_id else safe_name

    return f'''<!DOCTYPE html><html lang="zh"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1,maximum-scale=3.0,user-scalable=no">
<title>{filename}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;color:#e6edf3;min-height:100dvh;display:flex;flex-direction:column}}
.toolbar{{display:flex;align-items:center;justify-content:center;gap:12px;padding:10px;background:#161b22;border-bottom:1px solid #21262d;flex-shrink:0}}
.toolbar button{{background:#21262d;color:#c9d1d9;border:1px solid #30363d;border-radius:6px;padding:6px 14px;font-size:14px;cursor:pointer}}
.toolbar button:active{{background:#30363d}}
.toolbar button.dl{{background:#1a3a2a;border-color:#238636;color:#3fb950}}
.toolbar button.dl:active{{background:#2ea043}}
.toolbar span{{color:#8b949e;font-size:13px;min-width:56px;text-align:center}}
.canvas-area{{flex:1;overflow:auto;display:flex;justify-content:center;align-items:flex-start;padding:10px;-webkit-overflow-scrolling:touch}}
.canvas-wrap{{transform-origin:top center;transition:transform .15s}}
canvas{{display:block;box-shadow:0 2px 16px rgba(0,0,0,.4);margin:0 auto}}
.loading{{display:flex;align-items:center;justify-content:center;height:200px;color:#8b949e;font-size:14px}}
.error{{color:#f85149;text-align:center;padding:40px 20px}}
.error a{{color:#58a6ff}}
</style></head><body>
<div class="toolbar">
  <button onclick="prevPage()">◀</button>
  <span id="pageInfo">— / —</span>
  <button onclick="nextPage()">▶</button>
  <span id="zoomLabel" style="font-size:11px;min-width:auto">—</span>
  <button class="dl" onclick="downloadFile()">⬇ 下载</button>
</div>
<div class="canvas-area" id="canvasArea" ontouchstart="onTouchStart(event)" ontouchmove="onTouchMove(event)" ontouchend="onTouchEnd(event)">
  <div class="loading" id="loading">正在加载 PDF…</div>
</div>
<script src="/__openclaw__/canvas/pdf.min.js"></script>
<script>
pdfjsLib.GlobalWorkerOptions.workerSrc='/__openclaw__/canvas/pdf.worker.min.js';
var pdf=null,pageNum=1,total=0,scale=1.5,zoom=1,minZoom=0.5,maxZoom=5.0;
var pinchStart=0,pinchZoom=1;

function setZoom(z){{
  zoom=Math.max(minZoom,Math.min(maxZoom,z));
  var w=document.getElementById('canvasWrap');
  if(w)w.style.transform='scale('+zoom+')';
  document.getElementById('zoomLabel').textContent=Math.round(zoom*100)+'%';
}}

function onTouchStart(e){{
  if(e.touches&&e.touches.length==2){{
    pinchStart=Math.hypot(e.touches[0].clientX-e.touches[1].clientX,e.touches[0].clientY-e.touches[1].clientY);
    pinchZoom=zoom;
    e.preventDefault();
  }}
}}
function onTouchMove(e){{
  if(pinchStart&&e.touches&&e.touches.length==2){{
    var d=Math.hypot(e.touches[0].clientX-e.touches[1].clientX,e.touches[0].clientY-e.touches[1].clientY);
    if(Math.abs(d-pinchStart)>5){{ setZoom(pinchZoom*(d/pinchStart)); }}
    e.preventDefault();
  }}
}}
function onTouchEnd(e){{ pinchStart=0; }}

// 复用画廊下载机制：__HOST__ 由客户端替换为 Gateway 地址
window.__LITTLEHELPER_GALLERY__={{title:'{title if title != file_id else filename}',items:[
{{fileId:'{file_id}',displayName:'{disp}',fileName:'{safe_name}',
  downloadUrl:'{download_url}',
  mimeType:'application/pdf',thumbUrl:''}}
]}};

function renderPage(num,fitWidth){{
  pdf.getPage(num).then(function(page){{
    var vp=page.getViewport({{scale:scale}});
    // 自适应宽度：根据屏幕宽度计算初始 zoom
    if(fitWidth){{
      var areaW=document.getElementById('canvasArea').clientWidth-20;
      zoom=areaW/vp.width;
      zoom=Math.max(minZoom,Math.min(2.0,zoom));
    }}
    var canvas=document.createElement('canvas');
    var ctx=canvas.getContext('2d');
    canvas.height=vp.height;canvas.width=vp.width;
    var wrap=document.createElement('div');
    wrap.id='canvasWrap';wrap.className='canvas-wrap';
    wrap.style.transform='scale('+zoom+')';
    wrap.appendChild(canvas);
    var area=document.getElementById('canvasArea');
    area.innerHTML='';area.appendChild(wrap);
    page.render({{canvasContext:ctx,viewport:vp}});
    document.getElementById('pageInfo').textContent=num+' / '+total;
    document.getElementById('zoomLabel').textContent=Math.round(zoom*100)+'%';
    document.getElementById('loading').style.display='none';
  }});
}}

function prevPage(){{ if(pageNum>1){{ pageNum--;renderPage(pageNum); }} }}
function nextPage(){{ if(pageNum<total){{ pageNum++;renderPage(pageNum); }} }}
function downloadFile(){{ location.href='littlehelper://gallery/download?index=0'; }}

pdfjsLib.getDocument('{pdf_url}').promise.then(function(p){{
  pdf=p;total=p.numPages;pageNum=1;
  if(total===0){{ document.getElementById('canvasArea').innerHTML='<div class="error">PDF 无内容</div>';return; }}
  renderPage(1,true);
}}).catch(function(e){{
  document.getElementById('loading').style.display='none';
  document.getElementById('canvasArea').innerHTML='<div class="error">⚠️ PDF 加载失败<br><small>'+e.message+'</small><br><br><a href="{pdf_url}" target="_blank">📥 直接下载</a></div>';
}});
</script></body></html>'''


def viewer_image(filepath: str, filename: str) -> str:
    copy_to_canvas(filepath, filename)
    return f'''<!DOCTYPE html><html lang="zh"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1,maximum-scale=3.0,user-scalable=no">
<title>{filename}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;min-height:100dvh;display:flex;align-items:center;justify-content:center;padding:8px}}
img{{max-width:100%;max-height:100dvh;object-fit:contain;border-radius:4px}}
</style></head><body>
<img src="/__openclaw__/canvas/{filename}" alt="{filename}">
</body></html>'''


def viewer_docx(filepath: str, filename: str) -> str:
    """Word → 段落文字 HTML"""
    from docx import Document
    doc = Document(filepath)
    items = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if not t:
            items.append('<br>')
            continue
        sty = p.style.name if p.style else ''
        if 'Heading' in sty or 'Title' in sty or 'heading' in sty.lower():
            items.append(f'<h3>{t}</h3>')
        else:
            items.append(f'<p>{t}</p>')
    for table in doc.tables:
        items.append('<table>')
        for row in table.rows:
            cells = ''.join(f'<td>{c.text}</td>' for c in row.cells)
            items.append(f'<tr>{cells}</tr>')
        items.append('</table>')
    body = '\n'.join(items) if items else '<p style="color:#8b949e">(空文档)</p>'
    return _wrap_html(filename, '📄', body)


def viewer_xlsx(filepath: str, filename: str) -> str:
    """Excel → HTML 表格"""
    from openpyxl import load_workbook
    wb = load_workbook(filepath, data_only=True)
    parts = []
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f'<h3>📊 {name}</h3><table>')
        for row in ws.iter_rows(values_only=True):
            cells = ''.join(f'<td>{c if c is not None else ""}</td>' for c in row)
            parts.append(f'<tr>{cells}</tr>')
        parts.append('</table>')
    body = '\n'.join(parts) if parts else '<p style="color:#8b949e">(空表格)</p>'
    return _wrap_html(filename, '📊', body)


def viewer_pptx(filepath: str, filename: str) -> str:
    """PPT → 每页文字 + 表格"""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(f'<p>{t}</p>')
            if shape.has_table:
                texts.append('<table>')
                for row in shape.table.rows:
                    cells = ''.join(f'<td>{c.text}</td>' for c in row.cells)
                    texts.append(f'<tr>{cells}</tr>')
                texts.append('</table>')
        if texts:
            slides.append(f'<div class="slide"><h3>第 {i+1} 页</h3>{"".join(texts)}</div>')
    body = '\n'.join(slides) if slides else '<p style="color:#8b949e">(无文字内容)</p>'
    return _wrap_html(filename, '📽️', body)


def _wrap_html(filename: str, icon: str, body: str) -> str:
    return f'''<!DOCTYPE html><html lang="zh"><head><meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1,maximum-scale=1,user-scalable=no">
<title>{filename}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{background:#0d1117;color:#d2d8e0;padding:16px;font-family:-apple-system,'Microsoft YaHei',sans-serif;font-size:15px;line-height:1.8;min-height:100dvh}}
h3{{color:#f0f6fc;font-size:16px;margin:14px 0 6px}}
p{{margin:4px 0}}
table{{width:100%;border-collapse:collapse;margin:10px 0;font-size:13px}}
td{{border:1px solid #21262d;padding:6px 8px;vertical-align:top}}
.slide{{margin-bottom:20px;padding:12px;background:#161b22;border-radius:8px}}
.header{{color:#8b949e;font-size:12px;margin-bottom:16px;padding-bottom:8px;border-bottom:1px solid #21262d}}
</style></head><body>
<div class="header">{icon} {filename}</div>
{body}
</body></html>'''


def detect_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext == '.pdf': return 'pdf'
    if ext == '.docx': return 'docx'
    if ext == '.xlsx': return 'xlsx'
    if ext == '.pptx': return 'pptx'
    if ext in {'.txt','.md','.csv','.log','.json','.xml','.html','.htm','.css','.js','.py','.yaml','.yml','.ini','.cfg','.sh','.bat','.ps1','.sql'}: return 'txt'
    if ext in {'.jpg','.jpeg','.png','.gif','.webp','.bmp'}: return 'image'
    return 'unknown'


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({"error": "缺少 fileId"}, ensure_ascii=False))
        sys.exit(1)

    file_id = sys.argv[1]
    title = sys.argv[2] if len(sys.argv) > 2 else file_id

    resolved = resolve_file(file_id)
    if not resolved:
        print(json.dumps({"error": f"未找到文件: {file_id}"}, ensure_ascii=False))
        sys.exit(1)

    src_path, filename = resolved
    ftype = detect_type(filename)
    cleanup_old_viewers()

    if ftype == 'pdf':
        html = viewer_pdf(src_path, filename, file_id=file_id, title=title)
    elif ftype == 'txt':
        html = viewer_txt(src_path, filename)
    elif ftype == 'image':
        html = viewer_image(src_path, filename)
    elif ftype == 'docx':
        html = viewer_docx(src_path, filename)
    elif ftype == 'xlsx':
        html = viewer_xlsx(src_path, filename)
    elif ftype == 'pptx':
        html = viewer_pptx(src_path, filename)
    else:
        print(json.dumps({"error": f"不支持的文件格式: {filename}", "type": ftype}, ensure_ascii=False))
        sys.exit(1)

    viewer_filename = f'viewer_{file_id}.html'
    viewer_path = os.path.join(CANVAS_DIR, viewer_filename)
    with open(viewer_path, 'w', encoding='utf-8') as f:
        f.write(html)

    print(json.dumps({
        "ok": True,
        "type": ftype,
        "filename": filename,
        "viewer": viewer_filename,
        "url": f"/__openclaw__/canvas/{viewer_filename}",
        "title": title
    }, ensure_ascii=False))
